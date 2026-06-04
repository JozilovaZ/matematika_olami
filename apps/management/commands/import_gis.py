"""GIS masofaviy kursini saytga import qiladi.

Har bir hafta (1-15) -> bitta Dars.
  - NAZARIY.docx  -> dars.nazariy (HTML, rasmlar bilan)
  - TAQDIMOT/.pptx -> dars.taqdimot (slayd matni) + dars.taqdimot_fayl (yuklab olish)
  - AMALIY_USLUBIY_KO'RSATMA.docx -> dars.amaliy_uslubiy (HTML, rasmlar bilan)
  - TOPSHIRIQ.docx -> dars.topshiriq_matni (HTML, rasmlar bilan)
  - TEST.docx -> Topshiriq (test) + Savol lar
"""
import glob
import os
import re
import shutil
import zipfile
import random
import xml.etree.ElementTree as ET

import docx
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx.oxml.ns import qn
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table, _Cell
from docx.text.paragraph import Paragraph
from docx.document import Document as _DocClass

from django.conf import settings
from django.core.management.base import BaseCommand
from django.utils.html import escape

from apps.models import Kategoriya, Dars, Topshiriq, Savol

GIS_DIR = os.path.join(settings.BASE_DIR, 'GIS masofaviy')
KATEGORIYA_NOMI = 'Geoaxborot texnologiyalari (Masofaviy)'

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def iter_block_items(parent):
    """Hujjat tartibida paragraf va jadvallarni qaytaradi."""
    if isinstance(parent, _DocClass):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        parent_elm = parent._element
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


EMU_PER_PX = 9525.0


def _build_numbering(doc):
    """numId -> ilvl -> {'fmt', 'text'} xaritasini quradi."""
    mapping = {}
    try:
        npart = doc.part.numbering_part
    except (NotImplementedError, KeyError, AttributeError):
        return mapping
    el = npart.element
    num2abs = {}
    for n in el.findall(qn('w:num')):
        a = n.find(qn('w:abstractNumId'))
        if a is not None:
            num2abs[n.get(qn('w:numId'))] = a.get(qn('w:val'))
    abs_map = {}
    for ab in el.findall(qn('w:abstractNum')):
        lvls = {}
        for lvl in ab.findall(qn('w:lvl')):
            fmt = lvl.find(qn('w:numFmt'))
            txt = lvl.find(qn('w:lvlText'))
            lvls[lvl.get(qn('w:ilvl'))] = {
                'fmt': fmt.get(qn('w:val')) if fmt is not None else 'decimal',
                'text': txt.get(qn('w:val')) if txt is not None else '%1.',
            }
        abs_map[ab.get(qn('w:abstractNumId'))] = lvls
    for nid, aid in num2abs.items():
        mapping[nid] = abs_map.get(aid, {})
    return mapping


def _para_numinfo(para):
    """Paragrafning (numId, ilvl) yoki None."""
    pPr = para._p.pPr
    if pPr is None:
        return None
    npr = pPr.find(qn('w:numPr'))
    if npr is None:
        return None
    ni = npr.find(qn('w:numId'))
    if ni is None:
        return None
    il = npr.find(qn('w:ilvl'))
    return (ni.get(qn('w:val')), il.get(qn('w:val')) if il is not None else '0')


def _roman(n):
    vals = [(1000, 'm'), (900, 'cm'), (500, 'd'), (400, 'cd'), (100, 'c'),
            (90, 'xc'), (50, 'l'), (40, 'xl'), (10, 'x'), (9, 'ix'),
            (5, 'v'), (4, 'iv'), (1, 'i')]
    s = ''
    for v, sym in vals:
        while n >= v:
            s += sym
            n -= v
    return s


def _fmt_counter(n, fmt):
    if fmt == 'lowerLetter':
        return chr(96 + ((n - 1) % 26) + 1)
    if fmt == 'upperLetter':
        return chr(64 + ((n - 1) % 26) + 1)
    if fmt == 'lowerRoman':
        return _roman(n)
    if fmt == 'upperRoman':
        return _roman(n).upper()
    return str(n)


def _shd_fill(tcPr):
    if tcPr is None:
        return None
    shd = tcPr.find(qn('w:shd'))
    if shd is None:
        return None
    fill = shd.get(qn('w:fill'))
    if fill and fill != 'auto' and re.match(r'^[0-9A-Fa-f]{6}$', fill or ''):
        return '#' + fill
    return None


def _run_inner(run, doc, img_dir, img_url_base):
    """Run -> rasm(lar) + Word formatlashini saqlagan matn HTML."""
    out = ''
    el = run._element
    for drawing in el.findall(qn('w:drawing')):
        blip = drawing.find('.//' + qn('a:blip'))
        if blip is None:
            continue
        rid = blip.get(qn('r:embed'))
        if not rid:
            continue
        try:
            part = doc.part.related_parts[rid]
        except KeyError:
            continue
        ext = os.path.splitext(part.partname)[1] or '.png'
        fname = 'img_%s%s' % (rid, ext)
        os.makedirs(img_dir, exist_ok=True)
        with open(os.path.join(img_dir, fname), 'wb') as f:
            f.write(part.blob)
        ext_el = drawing.find('.//' + qn('wp:extent'))
        w_attr = ''
        if ext_el is not None and ext_el.get('cx'):
            w_attr = ' style="width:%.0fpx;"' % (int(ext_el.get('cx')) / EMU_PER_PX)
        out += '<img class="w-img"%s src="%s/%s" alt="">' % (w_attr, img_url_base, fname)
    text = run.text
    if text:
        text = escape(text).replace('\n', '<br>').replace('\t', '&emsp;')
        style = ''
        if run.font.size is not None:
            style += 'font-size:%gpt;' % run.font.size.pt
        if run.bold:
            style += 'font-weight:bold;'
        if run.italic:
            style += 'font-style:italic;'
        if run.underline:
            style += 'text-decoration:underline;'
        col = _font_color(run.font)
        if col:
            style += 'color:%s;' % col
        if run.font.name:
            style += "font-family:'%s';" % run.font.name
        out += ('<span style="%s">%s</span>' % (style, text)) if style else text
    return out


def _para_style(para):
    """Paragraf darajasidagi CSS (tekislash, chekinish, interval)."""
    pf = para.paragraph_format
    style = ''
    al = para.alignment
    if al is not None and int(al) in (1, 2, 3):
        style += 'text-align:%s;' % {1: 'center', 2: 'right', 3: 'justify'}[int(al)]
    if pf.first_line_indent:
        style += 'text-indent:%.0fpx;' % (pf.first_line_indent / EMU_PER_PX)
    if pf.left_indent:
        style += 'margin-left:%.0fpx;' % (pf.left_indent / EMU_PER_PX)
    if pf.space_before:
        style += 'margin-top:%.0fpx;' % (pf.space_before / EMU_PER_PX)
    if pf.space_after:
        style += 'margin-bottom:%.0fpx;' % (pf.space_after / EMU_PER_PX)
    return style


def _para_inner(para, doc, img_dir, img_url_base):
    inner = ''
    for run in para.runs:
        inner += _run_inner(run, doc, img_dir, img_url_base)
    return inner


def _table_html(table, doc, img_dir, img_url_base):
    rows_html = []
    for row in table.rows:
        cells_html = []
        for tc in row._tr.findall(qn('w:tc')):
            cell = _Cell(tc, table)
            tcPr = tc.find(qn('w:tcPr'))
            attr = ''
            cstyle = ''
            if tcPr is not None:
                span = tcPr.find(qn('w:gridSpan'))
                if span is not None:
                    attr = ' colspan="%s"' % span.get(qn('w:val'))
                fill = _shd_fill(tcPr)
                if fill:
                    cstyle += 'background:%s;' % fill
            content = ''
            for p in cell.paragraphs:
                pin = _para_inner(p, doc, img_dir, img_url_base)
                if pin.strip():
                    content += '<p style="%s">%s</p>' % (_para_style(p), pin)
            cells_html.append('<td%s style="%s">%s</td>' % (attr, cstyle, content or '&nbsp;'))
        rows_html.append('<tr>%s</tr>' % ''.join(cells_html))
    return '<table class="w-table">%s</table>' % ''.join(rows_html)


def docx_to_html(path, img_dir, img_url_base):
    """docx -> HTML, Word ko'rinishini imkon qadar saqlab (formatlash, ro'yxat,
    jadval, rasm, abzas joylashuvi)."""
    doc = docx.Document(path)
    numbering = _build_numbering(doc)
    counters = {}
    parts = []
    prev_numid = None
    prev_was_list = False

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            inner = _para_inner(block, doc, img_dir, img_url_base)
            if not inner.strip():
                prev_was_list = False
                prev_numid = None
                continue
            pstyle = _para_style(block)
            numinfo = _para_numinfo(block)

            if numinfo and numinfo[0] in numbering:
                numid, ilvl = numinfo
                lvl = numbering[numid].get(ilvl, {'fmt': 'decimal', 'text': '%1.'})
                # Yangi ro'yxat boshlanganda sanagichni qayta tiklash
                if not (prev_was_list and prev_numid == numid):
                    for k in list(counters):
                        if k[0] == numid:
                            del counters[k]
                ilvl_i = int(ilvl)
                if lvl['fmt'] == 'bullet':
                    marker = '◦' if ilvl_i > 0 else '•'
                else:
                    key = (numid, ilvl)
                    counters[key] = counters.get(key, 0) + 1
                    # chuqurroq darajalarni nolga tushirish
                    for k in list(counters):
                        if k[0] == numid and int(k[1]) > ilvl_i:
                            del counters[k]
                    marker = lvl['text']
                    for lv in range(9, -1, -1):
                        token = '%%%d' % (lv + 1)
                        if token in marker:
                            cnt = counters.get((numid, str(lv)), 1)
                            marker = marker.replace(token, _fmt_counter(cnt, lvl['fmt']))
                indent = 26 + ilvl_i * 26
                parts.append(
                    '<div class="w-li" style="margin-left:%dpx;%s">'
                    '<span class="w-marker">%s</span><span class="w-litext">%s</span></div>'
                    % (indent, pstyle, escape(marker), inner)
                )
                prev_was_list = True
                prev_numid = numid
            else:
                is_heading = bool(block.style) and (block.style.name or '').lower().startswith('heading')
                cls = ' class="w-h"' if is_heading else ''
                parts.append('<p%s style="%s">%s</p>' % (cls, pstyle, inner))
                prev_was_list = False
                prev_numid = None
        else:  # Table
            parts.append(_table_html(block, doc, img_dir, img_url_base))
            prev_was_list = False
            prev_numid = None

    return '<div class="w-doc">%s</div>' % ''.join(parts)


def _font_color(font):
    try:
        if font.color is not None and font.color.type is not None:
            return '#%s' % str(font.color.rgb)
    except Exception:
        pass
    return None


def _shape_fill(shape):
    try:
        fill = shape.fill
        if fill.type == 1:  # MSO_FILL.SOLID
            return '#%s' % str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _font_family(name):
    if not name:
        return None
    return 'Georgia, serif' if 'serif' in name.lower() else 'Arial, sans-serif'


def pptx_to_html(path, img_dir, img_url_base):
    """pptx slaydlarini PowerPointdagi joylashuvini saqlagan holda HTML qiladi.

    Har bir slayd 16:9 (yoki mos) kanvas; shakllar foizli koordinatalar bilan
    aniq joylashtiriladi, font o'lchami container-query birligida (cqw) beriladi.
    """
    try:
        prs = Presentation(path)
    except Exception:
        return ''
    sw = prs.slide_width or 12192000
    sh = prs.slide_height or 6858000
    sw_pt = sw / 12700.0  # slayd kengligi punktlarda (cqw hisobi uchun)

    def cqw(pt):
        return pt / sw_pt * 100.0

    slides_html = []
    for idx, slide in enumerate(prs.slides, 1):
        boxes = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
                continue
            pos = 'left:%.3f%%;top:%.3f%%;width:%.3f%%;height:%.3f%%;' % (
                shape.left / sw * 100, shape.top / sh * 100,
                shape.width / sw * 100, shape.height / sh * 100,
            )

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    fname = 'sl%d_%s.%s' % (idx, shape.shape_id, image.ext)
                    os.makedirs(img_dir, exist_ok=True)
                    with open(os.path.join(img_dir, fname), 'wb') as f:
                        f.write(image.blob)
                    boxes.append('<img class="pp-pic" style="%s" src="%s/%s" alt="">' % (
                        pos, img_url_base, fname))
                except Exception:
                    pass
                continue

            box_style = pos
            fill = _shape_fill(shape)
            if fill:
                box_style += 'background:%s;' % fill

            inner = ''
            if shape.has_text_frame and shape.text_frame.text.strip():
                tf = shape.text_frame
                va = 'center'
                try:
                    if tf.vertical_anchor == MSO_ANCHOR.TOP:
                        va = 'flex-start'
                    elif tf.vertical_anchor == MSO_ANCHOR.BOTTOM:
                        va = 'flex-end'
                except Exception:
                    pass
                paras = []
                for para in tf.paragraphs:
                    if not para.text.strip():
                        continue
                    align = 'left'
                    if para.alignment == PP_ALIGN.CENTER:
                        align = 'center'
                    elif para.alignment == PP_ALIGN.RIGHT:
                        align = 'right'
                    elif para.alignment == PP_ALIGN.JUSTIFY:
                        align = 'justify'
                    runs = []
                    for run in para.runs:
                        if not run.text:
                            continue
                        rs = ''
                        if run.font.size is not None:
                            rs += 'font-size:%.3fcqw;' % cqw(run.font.size.pt)
                        if run.font.bold:
                            rs += 'font-weight:700;'
                        if run.font.italic:
                            rs += 'font-style:italic;'
                        col = _font_color(run.font)
                        if col:
                            rs += 'color:%s;' % col
                        fam = _font_family(run.font.name)
                        if fam:
                            rs += 'font-family:%s;' % fam
                        runs.append('<span style="%s">%s</span>' % (rs, escape(run.text)))
                    if not runs:
                        runs.append(escape(para.text))
                    paras.append('<p style="text-align:%s;">%s</p>' % (align, ''.join(runs)))
                inner = ''.join(paras)
                box_style += 'display:flex;flex-direction:column;justify-content:%s;' % va

            if inner or fill:
                boxes.append('<div class="pp-box" style="%s">%s</div>' % (box_style, inner))

        slides_html.append(
            '<div class="pp-slide">'
            '<div class="pp-num">%d-slayd</div>'
            '<div class="pp-canvas" style="aspect-ratio:%d/%d;">%s</div>'
            '</div>' % (idx, sw, sh, ''.join(boxes))
        )
    return ''.join(slides_html)


def parse_test(path):
    """TEST.docx jadvalidan savollarni o'qiydi.

    Ustunlar: 0=№, 1=Savol, 2=To'g'ri javob, 3-5=Muqobil javoblar.
    """
    doc = docx.Document(path)
    if not doc.tables:
        return []
    table = doc.tables[0]
    savollar = []
    for row in table.rows[1:]:
        cells = [c.text.strip() for c in row.cells]
        if len(cells) < 6:
            continue
        savol = cells[1]
        togri = cells[2]
        notogri = [cells[3], cells[4], cells[5]]
        if not savol or not togri:
            continue
        variantlar = [togri] + notogri
        random.shuffle(variantlar)
        togri_idx = variantlar.index(togri)
        savollar.append({
            'savol': savol,
            'variantlar': variantlar,
            'togri': 'abcd'[togri_idx],
        })
    return savollar


class Command(BaseCommand):
    help = 'GIS masofaviy kursini saytga import qiladi'

    def add_arguments(self, parser):
        parser.add_argument('--tozalash', action='store_true',
                            help='Avval mavjud GIS darslarini o\'chiradi')

    def handle(self, *args, **opts):
        random.seed(42)
        if not os.path.isdir(GIS_DIR):
            self.stderr.write('GIS papkasi topilmadi: %s' % GIS_DIR)
            return

        kategoriya, _ = Kategoriya.objects.get_or_create(nomi=KATEGORIYA_NOMI)

        if opts.get('tozalash'):
            Dars.objects.filter(kategoriya=kategoriya).delete()
            self.stdout.write('Eski GIS darslari o\'chirildi.')

        media_root = str(settings.MEDIA_ROOT)
        taqdimot_dir = os.path.join(media_root, 'gis', 'taqdimotlar')
        os.makedirs(taqdimot_dir, exist_ok=True)

        for w in range(1, 16):
            wdir = os.path.join(GIS_DIR, '%d-hafta' % w)
            if not os.path.isdir(wdir):
                continue

            def find(pattern):
                res = glob.glob(os.path.join(wdir, pattern))
                return res[0] if res else None

            naz = find('*NAZARIY*.docx')
            amaliy = find('*AMALIY*.docx')
            topf = find('*TOPSHIRIQ*.docx')
            testf = find('*TEST*.docx')
            pptxf = find('*.pptx')

            # Mavzu nomi: AMALIY birinchi qatori, bo'lmasa NAZARIY
            nomi = '%d-mavzu' % w
            for src in (amaliy, naz):
                if src:
                    d = docx.Document(src)
                    for p in d.paragraphs:
                        if p.text.strip():
                            nomi = p.text.strip()
                            break
                    if nomi != '%d-mavzu' % w:
                        break

            def dirs_for(kind):
                d = os.path.join(media_root, 'gis', 'rasmlar', '%d-hafta' % w, kind)
                u = '%sgis/rasmlar/%d-hafta/%s' % (settings.MEDIA_URL, w, kind)
                return d, u

            dars, created = Dars.objects.update_or_create(
                kategoriya=kategoriya,
                raqam=str(w),
                defaults={
                    'nomi': nomi[:200],
                    'tavsif': nomi[:200],
                    'icon': '🌍',
                    'rang': '#e3f2fd',
                    'davomiylik': '4 soat',
                    'holat': 'tugallangan',
                    'tartib': w,
                },
            )

            if naz:
                dars.nazariy = docx_to_html(naz, *dirs_for('nazariy'))
            if amaliy:
                dars.amaliy_uslubiy = docx_to_html(amaliy, *dirs_for('amaliy'))
            if topf:
                dars.topshiriq_matni = docx_to_html(topf, *dirs_for('topshiriq'))
            if pptxf:
                dars.taqdimot = pptx_to_html(pptxf, *dirs_for('taqdimot'))
                dest_name = 'gis/taqdimotlar/%d-hafta%s' % (w, os.path.splitext(pptxf)[1])
                shutil.copy(pptxf, os.path.join(media_root, dest_name))
                dars.taqdimot_fayl.name = dest_name
            dars.save()

            # Test -> Topshiriq + Savollar
            if testf:
                savollar = parse_test(testf)
                if savollar:
                    top, _ = Topshiriq.objects.update_or_create(
                        dars=dars,
                        turi='kichik',
                        defaults={
                            'nomi': '%d-mavzu testi' % w,
                            'tavsif': '%s mavzusi bo\'yicha test' % nomi[:150],
                            'icon': '📝',
                            'rang': '#e8f5e9',
                            'savollar': len(savollar),
                            'tartib': w,
                        },
                    )
                    top.savollar_list.all().delete()
                    for i, s in enumerate(savollar):
                        Savol.objects.create(
                            topshiriq=top,
                            savol_matni=s['savol'],
                            variant_a=s['variantlar'][0][:200],
                            variant_b=s['variantlar'][1][:200],
                            variant_c=s['variantlar'][2][:200],
                            variant_d=s['variantlar'][3][:200],
                            togri_javob=s['togri'],
                            tartib=i,
                        )

            self.stdout.write(self.style.SUCCESS(
                '%d-hafta: %s%s' % (w, nomi[:60], '' if created else ' (yangilandi)')
            ))

        self.stdout.write(self.style.SUCCESS('\nImport tugadi.'))
